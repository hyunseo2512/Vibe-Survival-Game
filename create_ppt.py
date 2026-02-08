"""
Vibe Survival Game - Project Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color Palette
BG_DARK = RGBColor(0x11, 0x11, 0x11)
BG_SLIDE = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
ACCENT_GREEN = RGBColor(0x00, 0xFF, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x44, 0x44)
ACCENT_YELLOW = RGBColor(0xFF, 0xFF, 0x00)
ACCENT_BLUE = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GRAY = RGBColor(0x66, 0x66, 0x66)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if alpha is not None:
        from pptx.oxml.ns import qn
        solidFill = shape.fill._fill
        srgbClr = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha_elem = srgbClr.makeelement(qn('a:alpha'), {})
            alpha_elem.set('val', str(int(alpha * 1000)))
            srgbClr.append(alpha_elem)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT_GREEN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Add bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = 'Arial'

        p.space_after = Pt(6)

    return txBox


def add_colored_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_icon_box(slide, left, top, size, color, text, font_size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # ============================================================
    # SLIDE 1: Title Slide
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, BG_DARK)

    # Decorative top bar
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    # Player icon (green square)
    add_colored_rect(slide, Inches(5.9), Inches(1.8), Inches(0.6), Inches(0.6), ACCENT_GREEN)
    # Enemy icons
    add_colored_rect(slide, Inches(7.0), Inches(1.9), Inches(0.4), Inches(0.4), ACCENT_RED)
    add_colored_rect(slide, Inches(4.8), Inches(2.0), Inches(0.35), Inches(0.35), ACCENT_RED)
    # Bullet
    add_colored_rect(slide, Inches(6.6), Inches(2.05), Inches(0.15), Inches(0.15), ACCENT_YELLOW)

    # Title
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.2),
                 "VIBE SURVIVAL GAME", font_size=54, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(2), Inches(3.9), Inches(9.3), Inches(0.7),
                 "2D Top-Down Roguelike Survival Game", font_size=24,
                 color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, Inches(2.5), Inches(4.7), Inches(8.3), Inches(0.8),
                 "Phaser 3 + React + Electron  |  TypeScript  |  Desktop Application",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    add_colored_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    # ============================================================
    # SLIDE 2: Project Overview
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "PROJECT OVERVIEW", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    # Left card - What is it?
    add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(4.8), Inches(0.5),
                 "What is Vibe Survival?", font_size=22, color=ACCENT_GREEN, bold=True)

    add_bullet_text(slide, Inches(1.2), Inches(2.3), Inches(4.8), Inches(4),
                    [
                        "2D Top-Down Roguelike Survival Game",
                        "Minimalist auto-attack combat system",
                        "Endless wave survival with increasing difficulty",
                        "Character growth through leveling system",
                        "Desktop application built with Electron",
                        "Procedurally generated graphics (no external assets)",
                        "High score persistence across sessions",
                    ], font_size=15)

    # Right card - Key Numbers
    add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.3), BG_CARD)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(4.8), Inches(0.5),
                 "Key Numbers", font_size=22, color=ACCENT_YELLOW, bold=True)

    stats = [
        ("519", "Lines of Code"),
        ("7", "Source Files"),
        ("5", "Dependencies"),
        ("800x600", "Game Resolution"),
        ("32x32", "Sprite Size (px)"),
    ]
    y_pos = 2.4
    for val, label in stats:
        add_text_box(slide, Inches(7.4), Inches(y_pos), Inches(2), Inches(0.45),
                     val, font_size=28, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, Inches(9.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.4),
                     label, font_size=16, color=GRAY)
        y_pos += 0.85

    # ============================================================
    # SLIDE 3: Tech Stack
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "TECH STACK", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    tech_cards = [
        ("Game Engine", ACCENT_GREEN, [
            "Phaser 3 (v3.80.0)",
            "2D game framework",
            "Arcade physics engine",
            "Canvas renderer",
        ]),
        ("Frontend", ACCENT_BLUE, [
            "React 18 (v18.2.0)",
            "TypeScript (v5.3.3)",
            "JSX component architecture",
            "Phaser-React integration",
        ]),
        ("Desktop", ACCENT_PURPLE, [
            "Electron (v28.1.0)",
            "Cross-platform desktop app",
            "electron-store for data",
            "Preload script security",
        ]),
        ("Build Tools", ACCENT_YELLOW, [
            "Vite (v5.0.12)",
            "vite-plugin-electron",
            "electron-builder",
            "@vitejs/plugin-react",
        ]),
    ]

    x_pos = 0.6
    for title, color, items in tech_cards:
        card_w = Inches(2.9)
        add_shape_bg(slide, Inches(x_pos), Inches(1.5), card_w, Inches(5.2), BG_CARD)
        add_colored_rect(slide, Inches(x_pos), Inches(1.5), card_w, Inches(0.06), color)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(1.8), Inches(2.4), Inches(0.5),
                     title, font_size=20, color=color, bold=True)
        add_bullet_text(slide, Inches(x_pos + 0.3), Inches(2.5), Inches(2.4), Inches(3.5),
                        items, font_size=14, bullet_color=color)
        x_pos += 3.1

    # ============================================================
    # SLIDE 4: Architecture
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "ARCHITECTURE", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    # Architecture flow diagram using cards
    # Layer 1: Electron
    add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2), BG_CARD)
    add_colored_rect(slide, Inches(0.8), Inches(1.5), Inches(0.06), Inches(1.2), ACCENT_PURPLE)
    add_text_box(slide, Inches(1.2), Inches(1.6), Inches(3), Inches(0.4),
                 "Electron (Main Process)", font_size=18, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(0.5),
                 "main.ts  ->  BrowserWindow (800x600)  ->  preload.ts (electron-store IPC)",
                 font_size=14, color=GRAY)

    # Arrow
    add_text_box(slide, Inches(6), Inches(2.75), Inches(1.3), Inches(0.5),
                 "▼", font_size=24, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Layer 2: React
    add_shape_bg(slide, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.2), BG_CARD)
    add_colored_rect(slide, Inches(0.8), Inches(3.1), Inches(0.06), Inches(1.2), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), Inches(3.2), Inches(3), Inches(0.4),
                 "React (Renderer Process)", font_size=18, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(3.6), Inches(10), Inches(0.5),
                 "main.tsx  ->  App.tsx (Phaser Game initialization)  ->  <div id=\"phaser-container\">",
                 font_size=14, color=GRAY)

    # Arrow
    add_text_box(slide, Inches(6), Inches(4.35), Inches(1.3), Inches(0.5),
                 "▼", font_size=24, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Layer 3: Phaser Scenes
    add_shape_bg(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(2.2), BG_CARD)
    add_colored_rect(slide, Inches(0.8), Inches(4.7), Inches(0.06), Inches(2.2), ACCENT_GREEN)
    add_text_box(slide, Inches(1.2), Inches(4.8), Inches(4.5), Inches(0.4),
                 "StartScene.ts", font_size=18, color=ACCENT_GREEN, bold=True)
    add_bullet_text(slide, Inches(1.2), Inches(5.3), Inches(4.5), Inches(1.5),
                    ["Title screen & menu", "High score display", "Start button with hover effect"],
                    font_size=13, bullet_color=ACCENT_GREEN)

    add_shape_bg(slide, Inches(7.0), Inches(4.7), Inches(5.5), Inches(2.2), BG_CARD)
    add_colored_rect(slide, Inches(7.0), Inches(4.7), Inches(0.06), Inches(2.2), ACCENT_YELLOW)
    add_text_box(slide, Inches(7.4), Inches(4.8), Inches(4.5), Inches(0.4),
                 "GameScene.ts (288 lines)", font_size=18, color=ACCENT_YELLOW, bold=True)
    add_bullet_text(slide, Inches(7.4), Inches(5.3), Inches(4.5), Inches(1.5),
                    ["Main gameplay loop", "Physics & collision", "Auto-attack & leveling", "Enemy spawning system"],
                    font_size=13, bullet_color=ACCENT_YELLOW)

    # Scene transition arrow
    add_text_box(slide, Inches(6.1), Inches(5.5), Inches(1.2), Inches(0.5),
                 "◀ ▶", font_size=20, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 5: Core Game Mechanics
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "CORE GAME MECHANICS", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    # Player card
    add_shape_bg(slide, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), BG_CARD)
    add_colored_rect(slide, Inches(1.4), Inches(1.7), Inches(0.5), Inches(0.5), ACCENT_GREEN)
    add_text_box(slide, Inches(2.1), Inches(1.7), Inches(2), Inches(0.5),
                 "PLAYER", font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_text(slide, Inches(1.0), Inches(2.4), Inches(3.0), Inches(4),
                    [
                        "HP: 100",
                        "Speed: 200 px/s",
                        "Movement: WASD / Arrow",
                        "Auto-attack (range: 300px)",
                        "Fire rate: 2 shots/sec",
                        "Level-up: full heal",
                        "Scale +10% per level",
                    ], font_size=13, bullet_color=ACCENT_GREEN)

    # Combat card
    add_shape_bg(slide, Inches(4.8), Inches(1.4), Inches(3.8), Inches(5.5), BG_CARD)
    add_colored_rect(slide, Inches(5.6), Inches(1.7), Inches(0.5), Inches(0.5), ACCENT_YELLOW)
    add_text_box(slide, Inches(6.3), Inches(1.7), Inches(2), Inches(0.5),
                 "COMBAT", font_size=20, color=ACCENT_YELLOW, bold=True)
    add_bullet_text(slide, Inches(5.2), Inches(2.4), Inches(3.0), Inches(4),
                    [
                        "Auto-targeting system",
                        "Nearest enemy priority",
                        "Bullet speed: 400 px/s",
                        "Bullet lifetime: 1 sec",
                        "Damage: 5 per bullet",
                        "Knockback on collision",
                        "Red tint damage feedback",
                    ], font_size=13, bullet_color=ACCENT_YELLOW)

    # Enemy card
    add_shape_bg(slide, Inches(9.0), Inches(1.4), Inches(3.8), Inches(5.5), BG_CARD)
    add_colored_rect(slide, Inches(9.8), Inches(1.7), Inches(0.4), Inches(0.4), ACCENT_RED)
    add_text_box(slide, Inches(10.4), Inches(1.7), Inches(2), Inches(0.5),
                 "ENEMY", font_size=20, color=ACCENT_RED, bold=True)
    add_bullet_text(slide, Inches(9.4), Inches(2.4), Inches(3.0), Inches(4),
                    [
                        "HP: 10 (2 hits to kill)",
                        "Speed: 100 px/s",
                        "Spawn: every 1 second",
                        "Chase AI (moveToObject)",
                        "Contact damage: 10 HP",
                        "XP reward: 20 per kill",
                        "Random edge spawn",
                    ], font_size=13, bullet_color=ACCENT_RED)

    # ============================================================
    # SLIDE 6: Game State & Progression
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "GAME STATE & PROGRESSION", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    # State variables card
    add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(3.0), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(0.4),
                 "Game State Variables", font_size=20, color=ACCENT_BLUE, bold=True)

    state_code = (
        "hp: number = 100          // Current health\n"
        "maxHp: number = 100       // Maximum health\n"
        "xp: number = 0            // Experience points\n"
        "level: number = 1         // Current level\n"
        "score: number = 0         // Total score\n"
        "nextLevelXp: number = 100 // XP to next level\n"
        "isGameOver: boolean       // Game state flag"
    )
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5.2), Inches(2.2),
                 state_code, font_size=12, color=ACCENT_GREEN, font_name='Consolas')

    # Progression system card
    add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.0), BG_CARD)
    add_text_box(slide, Inches(7.4), Inches(1.7), Inches(5), Inches(0.4),
                 "Progression System", font_size=20, color=ACCENT_PURPLE, bold=True)
    add_bullet_text(slide, Inches(7.4), Inches(2.3), Inches(5), Inches(2.2),
                    [
                        "Kill enemies -> Earn 20 XP each",
                        "XP threshold increases +20% per level",
                        "Level up -> Full heal to max HP",
                        "Player sprite grows 10% per level",
                        "Score = XP earned x 10",
                        "High score saved to electron-store",
                    ], font_size=14, bullet_color=ACCENT_PURPLE)

    # Game flow
    add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(5.0), Inches(5), Inches(0.4),
                 "Game Flow", font_size=20, color=ACCENT_YELLOW, bold=True)

    flow_boxes = [
        ("START\nSCENE", ACCENT_BLUE),
        ("->", None),
        ("GAMEPLAY\n(Survive)", ACCENT_GREEN),
        ("->", None),
        ("ENEMIES\nSPAWN", ACCENT_RED),
        ("->", None),
        ("LEVEL UP\n(Heal)", ACCENT_PURPLE),
        ("->", None),
        ("GAME OVER\n(Score)", ACCENT_YELLOW),
        ("->", None),
        ("HIGH SCORE\n(Save)", ACCENT_BLUE),
    ]

    x = 1.0
    for text, color in flow_boxes:
        if color is None:
            add_text_box(slide, Inches(x), Inches(5.6), Inches(0.5), Inches(0.5),
                         text, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
            x += 0.5
        else:
            shape = add_shape_bg(slide, Inches(x), Inches(5.5), Inches(1.5), Inches(1.0), color)
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            p.font.color.rgb = BG_DARK
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            x += 1.7

    # ============================================================
    # SLIDE 7: Project Structure
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "PROJECT STRUCTURE", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    # File tree
    add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(0.4),
                 "Directory Layout", font_size=20, color=ACCENT_GREEN, bold=True)

    tree = (
        "Vibe-Survival-Game/\n"
        "  src/\n"
        "    main/\n"
        "      main.ts          (49 lines)\n"
        "      preload.ts       (3 lines)\n"
        "    renderer/\n"
        "      main.tsx         (11 lines)\n"
        "      App.tsx          (57 lines)\n"
        "      game/\n"
        "        GameScene.ts   (288 lines)\n"
        "        StartScene.ts  (60 lines)\n"
        "        Enemy.ts       (51 lines)\n"
        "  index.html\n"
        "  package.json\n"
        "  vite.config.ts\n"
        "  tsconfig.json"
    )
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5.2), Inches(4.5),
                 tree, font_size=13, color=GRAY, font_name='Consolas')

    # File descriptions
    add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD)
    add_text_box(slide, Inches(7.4), Inches(1.7), Inches(5), Inches(0.4),
                 "Key Files", font_size=20, color=ACCENT_YELLOW, bold=True)

    files = [
        ("main.ts", "Electron main process, window creation"),
        ("App.tsx", "React root, Phaser initialization"),
        ("GameScene.ts", "Core gameplay logic (288 lines)"),
        ("StartScene.ts", "Menu screen, high scores"),
        ("Enemy.ts", "Enemy AI and behavior"),
        ("vite.config.ts", "Build configuration"),
        ("package.json", "Dependencies & scripts"),
    ]

    y = 2.3
    for fname, desc in files:
        add_text_box(slide, Inches(7.4), Inches(y), Inches(2.5), Inches(0.35),
                     fname, font_size=14, color=ACCENT_BLUE, bold=True, font_name='Consolas')
        add_text_box(slide, Inches(9.6), Inches(y), Inches(2.8), Inches(0.35),
                     desc, font_size=13, color=GRAY)
        y += 0.6

    # ============================================================
    # SLIDE 8: Key Implementation Details
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "KEY IMPLEMENTATION DETAILS", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    # Procedural Graphics
    add_shape_bg(slide, Inches(0.6), Inches(1.4), Inches(5.9), Inches(2.6), BG_CARD)
    add_text_box(slide, Inches(1.0), Inches(1.6), Inches(5), Inches(0.4),
                 "Procedural Graphics", font_size=18, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(1.0), Inches(2.1), Inches(5.2), Inches(0.3),
                 "No external assets - all graphics generated at runtime",
                 font_size=13, color=GRAY)

    code_gfx = (
        "createTexture(key, color) {\n"
        "  const g = this.add.graphics();\n"
        "  g.fillStyle(color, 1);\n"
        "  g.fillRect(0, 0, 32, 32);\n"
        "  g.generateTexture(key, 32, 32);\n"
        "}"
    )
    add_text_box(slide, Inches(1.0), Inches(2.5), Inches(5.2), Inches(1.3),
                 code_gfx, font_size=11, color=ACCENT_GREEN, font_name='Consolas')

    # Auto-Attack System
    add_shape_bg(slide, Inches(6.9), Inches(1.4), Inches(5.9), Inches(2.6), BG_CARD)
    add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.4),
                 "Auto-Attack System", font_size=18, color=ACCENT_YELLOW, bold=True)
    add_text_box(slide, Inches(7.3), Inches(2.1), Inches(5.2), Inches(0.3),
                 "Automatic targeting of nearest enemy within range",
                 font_size=13, color=GRAY)

    code_atk = (
        "autoAttack() {\n"
        "  // Find nearest enemy within 300px\n"
        "  // Fire bullet at 400px/s velocity\n"
        "  // 0.5s cooldown between shots\n"
        "  // Bullet auto-destroys after 1s\n"
        "}"
    )
    add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5.2), Inches(1.3),
                 code_atk, font_size=11, color=ACCENT_YELLOW, font_name='Consolas')

    # Physics & Collision
    add_shape_bg(slide, Inches(0.6), Inches(4.3), Inches(5.9), Inches(2.6), BG_CARD)
    add_text_box(slide, Inches(1.0), Inches(4.5), Inches(5), Inches(0.4),
                 "Physics & Collision", font_size=18, color=ACCENT_BLUE, bold=True)
    add_bullet_text(slide, Inches(1.0), Inches(5.0), Inches(5.2), Inches(1.8),
                    [
                        "Phaser Arcade Physics (no gravity)",
                        "World bounds collision for all entities",
                        "Knockback on player-enemy collision",
                        "Bullet-enemy overlap detection",
                    ], font_size=13, bullet_color=ACCENT_BLUE)

    # Data Persistence
    add_shape_bg(slide, Inches(6.9), Inches(4.3), Inches(5.9), Inches(2.6), BG_CARD)
    add_text_box(slide, Inches(7.3), Inches(4.5), Inches(5), Inches(0.4),
                 "Data Persistence", font_size=18, color=ACCENT_PURPLE, bold=True)
    add_bullet_text(slide, Inches(7.3), Inches(5.0), Inches(5.2), Inches(1.8),
                    [
                        "electron-store for high score storage",
                        "JSON file in user config directory",
                        "store.get('highScore', 0) to retrieve",
                        "store.set('highScore', score) to save",
                    ], font_size=13, bullet_color=ACCENT_PURPLE)

    # ============================================================
    # SLIDE 9: Development & Build
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "DEVELOPMENT & BUILD", font_size=36, color=WHITE, bold=True)
    add_colored_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

    # Commands card
    add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.8), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(0.4),
                 "NPM Scripts", font_size=20, color=ACCENT_GREEN, bold=True)

    commands = [
        ("npm install", "Install all dependencies"),
        ("npm run dev", "Start dev server + Electron"),
        ("npm run build", "Build production application"),
        ("npm run preview", "Preview production build"),
    ]
    y = 2.3
    for cmd, desc in commands:
        add_text_box(slide, Inches(1.3), Inches(y), Inches(2.3), Inches(0.35),
                     cmd, font_size=14, color=ACCENT_GREEN, bold=True, font_name='Consolas')
        add_text_box(slide, Inches(3.6), Inches(y), Inches(2.8), Inches(0.35),
                     desc, font_size=14, color=GRAY)
        y += 0.5

    # Dependencies card
    add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.8), BG_CARD)
    add_text_box(slide, Inches(7.4), Inches(1.7), Inches(5), Inches(0.4),
                 "Dependencies", font_size=20, color=ACCENT_YELLOW, bold=True)

    deps = [
        ("phaser", "v3.80.0", "Game engine"),
        ("react", "v18.2.0", "UI framework"),
        ("react-dom", "v18.2.0", "React DOM renderer"),
        ("electron", "v28.1.0", "Desktop framework"),
        ("electron-store", "v8.1.0", "Data persistence"),
    ]
    y = 2.3
    for name, ver, desc in deps:
        add_text_box(slide, Inches(7.5), Inches(y), Inches(1.8), Inches(0.3),
                     name, font_size=13, color=ACCENT_BLUE, bold=True, font_name='Consolas')
        add_text_box(slide, Inches(9.3), Inches(y), Inches(1), Inches(0.3),
                     ver, font_size=12, color=DARK_GRAY, font_name='Consolas')
        add_text_box(slide, Inches(10.3), Inches(y), Inches(2), Inches(0.3),
                     desc, font_size=13, color=GRAY)
        y += 0.45

    # Linux compatibility
    add_shape_bg(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(4.8), Inches(10), Inches(0.4),
                 "Platform Compatibility Notes", font_size=20, color=ACCENT_RED, bold=True)
    add_bullet_text(slide, Inches(1.2), Inches(5.4), Inches(10), Inches(1.5),
                    [
                        "Hardware acceleration disabled for Linux compatibility (app.disableHardwareAcceleration())",
                        "Canvas renderer forced instead of WebGL for cross-platform support",
                        "Electron-builder configured for multi-platform packaging",
                        "TypeScript strict mode enabled for type safety across all modules",
                    ], font_size=14, bullet_color=ACCENT_RED)

    # ============================================================
    # SLIDE 10: Summary / Thank You
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_colored_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0),
                 "VIBE SURVIVAL GAME", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(2.7), Inches(9.3), Inches(0.6),
                 "Project Summary", font_size=28, color=ACCENT_GREEN,
                 alignment=PP_ALIGN.CENTER)

    # Summary grid
    summary_items = [
        ("2D Roguelike", "Top-down survival\nwith auto-attack"),
        ("Modern Stack", "Phaser 3 + React\n+ Electron + TS"),
        ("Lean Codebase", "519 lines across\n7 source files"),
        ("Zero Assets", "All graphics are\nprocedurally generated"),
    ]

    x = 1.2
    for title, desc in summary_items:
        add_shape_bg(slide, Inches(x), Inches(3.6), Inches(2.4), Inches(2.2), BG_CARD)
        add_text_box(slide, Inches(x + 0.2), Inches(3.8), Inches(2.0), Inches(0.4),
                     title, font_size=16, color=ACCENT_GREEN, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.2), Inches(4.4), Inches(2.0), Inches(1.2),
                     desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
        x += 2.7

    add_text_box(slide, Inches(2), Inches(6.3), Inches(9.3), Inches(0.6),
                 "THANK YOU", font_size=32, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_colored_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "Vibe_Survival_Game_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
